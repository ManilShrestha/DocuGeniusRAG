from abc import ABC, abstractmethod
import os

class DocLoader:
    def __init__(self, doc_path) -> None:
        self.doc_path = doc_path
        self.doc_text = ""
        if not os.path.exists(self.doc_path):
            raise FileNotFoundError(f"File {self.doc_path} does not exist.")
        
        # Determine the correct loader based on the file extension
        _, ext = os.path.splitext(doc_path)
        if ext.lower() == '.pdf':
            self.loader = PDFDocLoader(doc_path)
        elif ext.lower() == '.docx':
            self.loader = DOCXDocLoader(doc_path)
        elif ext.lower() =='.pptx':
            self.loader = PPTXDocLoader(doc_path)
        else:
            raise NotImplementedError(f"No loader implemented for {ext} files.")
        
        self.doc_text = self.loader.read_contents()
    
    def chunkify_document(self, chunk_size=1000):
        """
        This method will chunkify the document into chunks based on a fixed number of characters,
        with a 50% overlap between consecutive chunks.

        Args:
            chunk_size (int, optional): The number of characters per chunk. Defaults to 1000.
        """
        # Initialize variables
        chunks = []
        start = 0
        half_chunk = int(chunk_size / 2)

        # Loop through the document and create chunks with 50% overlap
        while start < len(self.doc_text):
            # Define the end of the current chunk
            end = start + chunk_size
            
            # Ensure the end does not exceed the document length
            if end > len(self.doc_text):
                end = len(self.doc_text)
            
            # Append the current chunk to the chunks list
            chunks.append(self.doc_text[start:end])
            
            # Update the start for the next chunk to achieve 50% overlap
            start += half_chunk
            
            # Break the loop if the next start plus the chunk size exceeds the document length
            if start + chunk_size > len(self.doc_text):
                # Check if there is any text left to consider as a final chunk
                if start < len(self.doc_text):
                    chunks.append(self.doc_text[start:])
                break

        return chunks

    
    def get_doc_text(self):
        return self.doc_text


class MasterDocumentLoader(ABC):
    def __init__(self, doc_path) -> None:
        self.doc_path = doc_path
        self.doc_text = ""

    @abstractmethod
    def read_contents(self):
        pass

class PDFDocLoader(MasterDocumentLoader):
    def read_contents(self):
        # from pdfminer.high_level import extract_text
        # self.doc_text = extract_text(self.doc_path)
        # return self.doc_text
        # Open the provided PDF file

        import fitz
        document = fitz.open(self.doc_path)
        
        # Initialize a variable to store all the extracted text
        all_text = ""
        
        # Loop through each page in the document
        for page in document:
            # Extract text from the current page
            text = page.get_text() # type: ignore
            all_text += text
        
        # Close the document
        document.close()
        
        return all_text

class DOCXDocLoader(MasterDocumentLoader):
    def read_contents(self):
        import docx
        doc = docx.Document(self.doc_path)
        self.doc_text = " ".join([para.text for para in doc.paragraphs])
        # Handle table text similarly as shown previously
        return self.doc_text

class PPTXDocLoader(MasterDocumentLoader):
    def read_contents(self):
        import pptx
        prs = pptx.Presentation(self.doc_path)
        self.doc_text = " ".join([run.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs])
        return self.doc_text
